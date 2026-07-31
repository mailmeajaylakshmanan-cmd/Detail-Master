"""Generate Detailing Masters Billing architecture PPT."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

GOLD = RGBColor(0xFB, 0xD9, 0x04)
DARK = RGBColor(0x11, 0x18, 0x27)
BLUE = RGBColor(0x25, 0x63, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x64, 0x74, 0x8B)
LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
GREEN = RGBColor(0x05, 0x96, 0x69)
ROSE = RGBColor(0xE1, 0x1D, 0x48)


def set_run(run, size=14, bold=False, color=DARK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_bg(slide, color=DARK):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # send to back
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_gold_bar(slide, top=0):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(top), Inches(13.333), Inches(0.12)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()


def title_box(slide, text, left=0.5, top=0.35, width=12, height=0.6, size=28, color=DARK):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=True, color=color)
    return box


def body_box(slide, lines, left=0.5, top=1.2, width=12, height=5.5, size=16, color=DARK, bold_first=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = line
        set_run(run, size=size, bold=(bold_first and i == 0), color=color)
    return box


def card(slide, left, top, width, height, title, bullets, fill=LIGHT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    shape.adjustments[0] = 0.1

    t = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(width - 0.4), Inches(0.4))
    run = t.text_frame.paragraphs[0].add_run()
    run.text = title
    set_run(run, size=14, bold=True, color=BLUE)

    b = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.55), Inches(width - 0.4), Inches(height - 0.7))
    tf = b.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = "• " + line
        set_run(run, size=12, color=DARK)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 Title
    s = prs.slides.add_slide(blank)
    add_bg(s, DARK)
    add_gold_bar(s, 0)
    add_gold_bar(s, 7.38)
    title_box(s, "DETAILING MASTERS", 0.8, 2.2, 11.5, 0.7, 40, GOLD)
    title_box(s, "Billing App — Full Structure & UI Model", 0.8, 3.0, 11.5, 0.5, 24, WHITE)
    body_box(
        s,
        ["React + Express + Neon Postgres", "Architecture · Permissions · Invoice ID flow · Connected UI"],
        0.8, 3.8, 11, 1.5, 16, RGBColor(0xCB, 0xD5, 0xE1),
    )

    # 2 Agenda
    s = prs.slides.add_slide(blank)
    add_gold_bar(s)
    title_box(s, "Agenda", 0.5, 0.35)
    items = [
        "1. Stack & hosting",
        "2. UI shell (sidebar + header)",
        "3. Screens & routes",
        "4. Permission / menu ID flow",
        "5. Billing data ID flow",
        "6. Auth & sessions",
        "7. Connected vs not connected",
        "8. Create invoice end-to-end",
    ]
    body_box(s, items, 0.7, 1.2, 11, 5.5, 20)

    # 3 Stack
    s = prs.slides.add_slide(blank)
    add_gold_bar(s)
    title_box(s, "1. Stack & Hosting", 0.5, 0.35)
    card(s, 0.5, 1.3, 3.8, 4.8, "Frontend", [
        "React + Vite",
        "React Query caching",
        "Sidebar menus from login",
        "Local: localhost:5173",
        "Prod: Vercel (manage.detailingmasters.in)",
    ])
    card(s, 4.7, 1.3, 3.8, 4.8, "API", [
        "Express (Node)",
        "JWT + login_sessions",
        "Local: :4000",
        "Prod: Railway",
        "/api/* routes",
    ])
    card(s, 8.9, 1.3, 3.9, 4.8, "Database", [
        "Neon Postgres",
        "Region: AWS US East 2",
        "Pooled DATABASE_URL",
        "Tables: clients, invoices,",
        "payments, menus, roles…",
    ])

    # 4 UI shell
    s = prs.slides.add_slide(blank)
    add_gold_bar(s)
    title_box(s, "2. UI Shell Model", 0.5, 0.35)
    body_box(
        s,
        [
            "┌── Sidebar ──────────────┬── Header (Search · Bell · Add New) ──┐",
            "│ Logo                    │                                       │",
            "│ Dynamic menus           │           <Page Outlet />             │",
            "│ (from login user.menus) │                                       │",
            "│ Sign Out                │                                       │",
            "└─────────────────────────┴── Footer ─────────────────────────────┘",
            "",
            "Glass car background is global (body CSS) on all pages.",
            "Menus are NOT hardcoded — driven by roles / user_menus.",
        ],
        0.6, 1.15, 12, 5.8, 16,
    )

    # 5 Screens
    s = prs.slides.add_slide(blank)
    add_gold_bar(s)
    title_box(s, "3. Screens & Routes", 0.5, 0.35)
    card(s, 0.4, 1.15, 4.1, 5.6, "Connected screens", [
        "/login — Auth",
        "/ — Dashboard",
        "/master-customer — Clients + vehicles",
        "/master-service — Services",
        "/invoices — Invoice list",
        "/invoices/new — Create bill",
        "/invoices/:id — View + Add Payment",
        "/invoices/:id/edit — Edit",
        "/menu-assignment — Role menus",
        "/user-menu-assignment — User menus",
    ])
    card(s, 4.7, 1.15, 4.1, 5.6, "Sidebar tree (menus table)", [
        "Dashboard → /",
        "Masters",
        "  · Customers",
        "  · Offers",
        "  · Services",
        "Billing & Records → /invoices",
        "Web Bookings",
        "Permissions",
        "  · Role Menus",
        "  · User Menus",
    ])
    card(s, 9.0, 1.15, 3.8, 5.6, "Not connected yet", [
        "Offers master",
        "Assign / view offer",
        "Website bookings",
        "Update credentials",
        "(UI exists, no Neon table)",
    ])

    # 6 Permissions flow
    s = prs.slides.add_slide(blank)
    add_gold_bar(s)
    title_box(s, "4. Permission ID Flow", 0.5, 0.35)
    body_box(
        s,
        [
            "admin_users.id ──role_id──► roles.id",
            "roles.id ──role_id──► role_menus ◄──menu_id── menus.id",
            "admin_users.id ──user_id──► user_menus ◄──menu_id── menus.id",
            "menus.id ──parent_id──► menus.id  (tree / groups)",
            "",
            "Login:",
            "  Super Admin → all menus",
            "  Other roles → role_menus.can_view OR user_menus.can_view",
            "  Build tree → localStorage.user.menus → Sidebar",
            "",
            "Pages:",
            "  Role Menus  → POST /permissions/roles/:roleId/menus",
            "  User Menus  → POST /permissions/users/:userId/menus",
        ],
        0.6, 1.15, 12, 5.8, 16,
    )

    # 7 Billing flow
    s = prs.slides.add_slide(blank)
    add_gold_bar(s)
    title_box(s, "5. Billing Data ID Flow", 0.5, 0.35)
    body_box(
        s,
        [
            "clients.id ──client_id──► vehicles.id",
            "clients.id ──client_id──► invoices.id",
            "vehicles.id ──vehicle_id──► invoices.id",
            "invoices.id ──invoice_order_id──► invoice_services",
            "services.id ──service_id──► invoice_services",
            "invoices.id ──invoice_order_id──► payments",
            "",
            "Money formulas:",
            "  unit_price     = snapshot of services.base_price at save time",
            "  sub_total      = SUM(invoice_services.unit_price)",
            "  grand_total    = sub_total − discount",
            "  amount_paid    = SUM(payments.amount)",
            "  balance_due    = grand_total − amount_paid",
            "",
            "Invoice number format: INV-DM-{timestamp}",
        ],
        0.6, 1.15, 12, 5.8, 15,
    )

    # 8 Auth
    s = prs.slides.add_slide(blank)
    add_gold_bar(s)
    title_box(s, "6. Auth & Session", 0.5, 0.35)
    card(s, 0.5, 1.3, 6, 5.2, "Login", [
        "POST /auth/login",
        "Check admin_users password_hash",
        "Create JWT (1 day)",
        "Insert login_sessions (user_id, token)",
        "Return user + menus tree + token",
        "Frontend stores token + user in localStorage",
    ])
    card(s, 6.8, 1.3, 5.9, 5.2, "Protect middleware", [
        "Read cookie token OR Bearer header",
        "Verify JWT",
        "Check login_sessions (not logged out)",
        "Short session cache (30s)",
        "Used by /permissions/* APIs",
        "Logout sets logout_at on session",
    ])

    # 9 Tables
    s = prs.slides.add_slide(blank)
    add_gold_bar(s)
    title_box(s, "7. Neon Tables Map", 0.5, 0.35)
    card(s, 0.4, 1.2, 4.1, 5.4, "Auth / ACL", [
        "admin_users",
        "roles",
        "menus",
        "role_menus",
        "user_menus",
        "login_sessions",
    ])
    card(s, 4.7, 1.2, 4.1, 5.4, "Masters", [
        "clients",
        "vehicles (client_id)",
        "services (base_price)",
    ])
    card(s, 9.0, 1.2, 3.8, 5.4, "Billing", [
        "invoices",
        "invoice_services",
        "payments",
        "",
        "FK key field:",
        "invoice_order_id",
    ])

    # 10 Create invoice
    s = prs.slides.add_slide(blank)
    add_gold_bar(s)
    title_box(s, "8. Create Invoice — End to End", 0.5, 0.35)
    body_box(
        s,
        [
            "UI: New Invoice form",
            "  1. Select client  → client_id",
            "  2. Select vehicle → vehicle_id (from that client)",
            "  3. Select services → service_ids[] (price locked from services)",
            "  4. Optional payments → amount, method, date (+ ref for UPI/bank)",
            "",
            "API: one POST /api/invoices",
            "  → insert invoices (INV-DM-…)",
            "  → insert invoice_services (invoice_order_id = invoice.id)",
            "  → insert payments (invoice_order_id = invoice.id)",
            "  → recalculate amount_paid + balance_due",
            "",
            "Later: Invoice View → Add Payment → another payments row, same invoice id",
        ],
        0.6, 1.15, 12, 5.8, 15,
    )

    # 11 Speed notes
    s = prs.slides.add_slide(blank)
    add_gold_bar(s)
    title_box(s, "Performance Notes", 0.5, 0.35)
    body_box(
        s,
        [
            "Fastest practical setup for this project:",
            "  • Same region for Railway API + Neon DB",
            "  • Neon pooled connection string",
            "  • Paginated invoice list + React Query cache",
            "  • Lean queries (no over-fetch)",
            "",
            "Current known gap (from earlier check):",
            "  Neon = US East 2 (Ohio)",
            "  Railway was US West (California) → move Railway to US East",
            "",
            "Do hosting region setup AFTER schema/UI connections are solid.",
        ],
        0.6, 1.2, 12, 5.5, 16,
    )

    # 12 Closing
    s = prs.slides.add_slide(blank)
    add_bg(s, DARK)
    add_gold_bar(s, 0)
    add_gold_bar(s, 7.38)
    title_box(s, "Thank you", 0.8, 2.6, 11.5, 0.7, 40, GOLD)
    title_box(s, "Detailing Masters Billing — Architecture Overview", 0.8, 3.5, 11.5, 0.5, 18, WHITE)
    body_box(
        s,
        ["Questions? Next: connect Offers / Bookings, or finish Railway region setup."],
        0.8, 4.3, 11, 1, 14, RGBColor(0xCB, 0xD5, 0xE1),
    )

    out = r"c:\Our-project\detailing-masters-billing\Detailing-Masters-App-Structure.pptx"
    prs.save(out)
    print(out)


if __name__ == "__main__":
    build()
